import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from datetime import datetime, timedelta

# Load event log data
eventlog_df1 = pd.read_excel("Hardware_Eventlog.xlsx", sheet_name=None, engine='openpyxl')
eventlog_df = pd.concat(eventlog_df1.values(), keys=eventlog_df1.keys()).reset_index(level=0).rename(columns={'level_0': 'Server'})
eventlog_df['Last Update'] = pd.to_datetime(eventlog_df['Last Update'], errors='coerce')

# Generate simulated metric data
servers = eventlog_df['Server'].unique()[:5]
timestamps = pd.date_range(end=datetime.now(), periods=48, freq='H')
metric_data = []
for server in servers:
    for ts in timestamps:
        cpu = max(0, min(100, round(np.random.normal(50, 10), 2)))
        mem = max(0, min(100, round(np.random.normal(60, 15), 2)))
        disk = max(0, round(np.random.normal(100, 20), 2))
        net = max(0, round(np.random.normal(200, 50), 2))
        metric_data.append([server, ts, cpu, mem, disk, net])
metric_df = pd.DataFrame(metric_data, columns=['Server', 'Timestamp', 'CPU', 'Memory', 'Disk_IO', 'Network'])

# Detect anomalies
def detect_anomalies(df):
    anomalies = []
    for server in df['Server'].unique():
        sub = df[df['Server'] == server]
        for metric in ['CPU', 'Memory', 'Disk_IO', 'Network']:
            mean = sub[metric].mean()
            std = sub[metric].std()
            threshold = mean + 2 * std
            outliers = sub[sub[metric] > threshold]
            for _, row in outliers.iterrows():
                anomalies.append({
                    'Server': server,
                    'Timestamp': row['Timestamp'],
                    'Metric': metric,
                    'Value': row[metric],
                    'Deviation': row[metric] - mean,
                    'Threshold': threshold
                })
    return pd.DataFrame(anomalies)

anomaly_df = detect_anomalies(metric_df)

# Create a presentation
prs = Presentation()
bullet_slide_layout = prs.slide_layouts[1]

def add_slide(title, bullets, image_path=None):
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    shapes.title.text = title
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
    if image_path:
        slide.shapes.add_picture(image_path, Inches(5.5), Inches(1.5), width=Inches(4.5))

# Slide 1: Overview
add_slide("Solution Overview", [
    "Proactive Server/Device Failure Detection Dashboard",
    "Real-time monitoring, anomaly detection, prediction, and alerting",
    "Covers all MVP and advanced requirements"
])

# Slide 2: Data Ingestion
add_slide("Data Ingestion & Asset Scope", [
    "Ingests event logs, server issue data, and ticket dumps",
    "Simulates metrics (CPU, Memory, Disk I/O, Network) for 5+ servers",
    "Supports real and synthetic data"
])

# Slide 3: Baseline Modeling
fig, ax = plt.subplots()
sample = metric_df[metric_df['Server'] == servers[0]]
ax.plot(sample['Timestamp'], sample['CPU'], label='CPU')
ax.axhline(sample['CPU'].mean(), color='green', linestyle='--', label='Baseline')
ax.axhline(sample['CPU'].mean() + 2 * sample['CPU'].std(), color='red', linestyle='--', label='Anomaly Threshold')
ax.set_title('Baseline Modeling & Anomaly Detection')
ax.legend()
img3 = BytesIO()
plt.savefig(img3, format='png')
img3.seek(0)
with open("slide3.png", "wb") as f:
    f.write(img3.read())
add_slide("Baseline Modeling & Anomaly Detection", [
    "Establishes statistical baselines for each metric per server",
    "Detects anomalies using thresholds",
    "Highlights deviations from normal behavior"
], "slide3.png")

# Slide 4: TTF Prediction
add_slide("Time-to-Failure (TTF) Prediction", [
    "Predicts TTF for hardware and metrics using historical patterns",
    "Displays TTF estimates for each server and metric",
    "Enables proactive maintenance"
])

# Slide 5: Cross-Metric Analysis
heatmap_data = anomaly_df.pivot_table(index='Server', columns='Metric', values='Deviation', aggfunc='mean').fillna(0)
fig, ax = plt.subplots(figsize=(6, 4))
sns.heatmap(heatmap_data, annot=True, cmap='coolwarm', ax=ax)
plt.title("Cross-Metric Anomaly Heatmap")
img5 = BytesIO()
plt.savefig(img5, format='png')
img5.seek(0)
with open("slide5.png", "wb") as f:
    f.write(img5.read())
add_slide("Cross-Metric & Cross-Device Analysis", [
    "Visualizes anomalies across metrics and servers",
    "Detects cascading failures and correlated events",
    "Heatmaps and tables for root cause analysis"
], "slide5.png")

# Slide 6: Alerting
add_slide("Alerting & Notification", [
    "Sends alerts via Outlook email and Microsoft Teams",
    "Includes server, metric, deviation, timestamp",
    "Supports routing, deduplication, and history"
])

# Slide 7: Interactive Visualization
add_slide("Interactive Visualization & Customization", [
    "Interactive dashboards with Plotly and Seaborn",
    "Custom thresholds and alert settings",
    "Multi-server comparison and forecasts"
])

# Slide 8: User Experience & Roles
add_slide("User Experience & Roles", [
    "Role-based views: Viewer, Operator, Admin",
    "Export/download anomaly data",
    "Simulated remediation playbooks"
])

# Slide 9: MVP & Beyond
add_slide("MVP & Beyond", [
    "MVP fully covered: ingestion, detection, TTF, alerting",
    "Advanced features: thresholds, plots, history, heatmaps",
    "Scalable and extensible for AI/ML"
])

# Slide 10: Business Value
add_slide("Business Value", [
    "Enables proactive, data-driven IT operations",
    "Reduces downtime and improves reliability",
    "Optimizes maintenance and enhances customer satisfaction"
])

# Save presentation
pptx_filename = "Proactive_Failure_Detection_Dashboard_Presentation.pptx"
prs.save(pptx_filename)
print(f"Presentation saved as {pptx_filename}")
